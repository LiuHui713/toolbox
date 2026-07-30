"""
Toolbox - 工具箱网站
"""

import io, os, uuid, threading
from datetime import datetime
from flask import Flask, render_template, request, send_file, make_response, jsonify
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(__file__), 'uploads')
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf','pptx','ppt','docx','doc','png','jpg','jpeg','mp3','wav','mp4','avi','mov','webm','m4a','ogg','txt','md','csv'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Preload Whisper model at startup (lazy, cached)
_whisper_model = None
def get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        from faster_whisper import WhisperModel
        _whisper_model = WhisperModel("small", device="cpu", compute_type="int8")
    return _whisper_model

# Start model download in background
threading.Thread(target=get_whisper_model, daemon=True).start()


@app.route('/')
def home():
    return render_template('home.html')


@app.route('/text')
def text_tool():
    return render_template('text.html')

@app.route('/text/extract', methods=['POST'])
def text_extract():
    try:
        texts = []
        original_names = []
        files = request.files.getlist('files')
        
        for f in files:
            if not f.filename or not allowed_file(f.filename):
                continue
            ext = f.filename.rsplit('.', 1)[1].lower()
            fname = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(f.filename))
            f.save(fname)
            original_names.append(f.filename)

            try:
                if ext == 'pdf':
                    import fitz
                    doc = fitz.open(fname)
                    for page in doc:
                        texts.append(page.get_text())
                    doc.close()
                elif ext in ('pptx', 'ppt'):
                    from pptx import Presentation
                    prs = Presentation(fname)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                texts.append(shape.text_frame.text)
                elif ext in ('docx', 'doc'):
                    import docx as _docx
                    doc = _docx.Document(fname)
                    texts.extend([p.text for p in doc.paragraphs])
                elif ext in ('png', 'jpg', 'jpeg'):
                    from rapidocr_onnxruntime import RapidOCR
                    ocr = RapidOCR()
                    result, _ = ocr(fname)
                    if result:
                        texts.extend([line[1] for line in result])
                elif ext in ('mp3', 'wav', 'm4a', 'ogg'):
                    import requests as req
                    groq_key = os.environ.get('GROQ_API_KEY', '')
                    if groq_key:
                        with open(fname, 'rb') as af:
                            resp = req.post(
                                'https://api.groq.com/openai/v1/audio/transcriptions',
                                headers={'Authorization': f'Bearer {groq_key}'},
                                files={'file': (f.filename, af)},
                                data={'model': 'whisper-large-v3', 'language': 'zh'}
                            )
                        if resp.status_code == 200:
                            texts.append(resp.json()['text'])
                        else:
                            raise Exception(f"Groq API error: {resp.text}")
                    else:
                        model = get_whisper_model()
                        segments, _ = model.transcribe(fname, language="zh")
                        texts.extend([s.text for s in segments])
                elif ext in ('mp4', 'avi', 'mov', 'webm'):
                    import subprocess, requests as req
                    audio_path = fname + '.wav'
                    subprocess.run(['ffmpeg', '-i', fname, '-vn', '-acodec', 'pcm_s16le', '-y', audio_path],
                                   capture_output=True, timeout=120)
                    groq_key = os.environ.get('GROQ_API_KEY', '')
                    if groq_key:
                        with open(audio_path, 'rb') as af:
                            resp = req.post(
                                'https://api.groq.com/openai/v1/audio/transcriptions',
                                headers={'Authorization': f'Bearer {groq_key}'},
                                files={'file': (os.path.basename(fname)+'.wav', af)},
                                data={'model': 'whisper-large-v3', 'language': 'zh'}
                            )
                        if resp.status_code == 200:
                            texts.append(resp.json()['text'])
                        else:
                            raise Exception(f"Groq API error: {resp.text}")
                    else:
                        model = get_whisper_model()
                        segments, _ = model.transcribe(audio_path, language="zh")
                        texts.extend([s.text for s in segments])
                    os.remove(audio_path)
                elif ext in ('txt', 'md', 'csv'):
                    with open(fname, 'r', encoding='utf-8', errors='ignore') as fh:
                        texts.append(fh.read())
            except Exception as e:
                texts.append(f"[提取失败: {f.filename} - {str(e)}]")
            finally:
                os.remove(fname)

        if not texts:
            return jsonify({"error": "未能提取到文字"}), 400

        from docx import Document as DocxDoc
        doc = DocxDoc()
        doc.add_heading('文字提取结果', 0)
        for i, t in enumerate(texts):
            if t.strip():
                doc.add_paragraph(t.strip())
        
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        base = original_names[0].rsplit('.', 1)[0] if original_names else 'extracted'
        dl_name = f'{base}_{ts}.docx'
        
        from urllib.parse import quote
        resp = make_response(buf.read())
        resp.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        resp.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{quote(dl_name)}"
        resp.headers['Content-Length'] = str(len(resp.data))
        resp.headers['X-Content-Type-Options'] = 'nosniff'
        resp.headers['Cache-Control'] = 'no-cache'
        return resp

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/ppt')
def ppt_tool():
    return render_template('ppt.html')

@app.route('/ppt/generate', methods=['POST'])
def ppt_generate():
    try:
        content = request.form.get('content', '')
        title = request.form.get('title', '演示文稿')
        if not content.strip():
            return jsonify({"error": "请输入内容"}), 400

        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text = "AI 自动生成"

        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        for para in paragraphs:
            lines = para.split('\n')
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = lines[0][:80]
            body = slide.placeholders[1].text_frame
            for line in lines[1:8]:
                p = body.add_paragraph()
                p.text = line[:200]
                p.level = 0

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.document',
                         as_attachment=True, download_name='generated.pptx')
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/audio')
def audio_tool():
    return render_template('audio.html')


@app.route('/poster')
def poster_tool():
    return render_template('poster.html')

@app.route('/poster/generate', methods=['POST'])
def poster_generate():
    try:
        text = request.form.get('text', '')
        bg_color = request.form.get('bg_color', '#1a1a2e')
        text_color = request.form.get('text_color', '#ffffff')

        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (800, 1200), bg_color)
        draw = ImageDraw.Draw(img)

        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()

        lines = text.split('\n')
        y = 200
        for i, line in enumerate(lines):
            font = font_large if i == 0 else font_small
            bbox = draw.textbbox((0, 0), line, font=font)
            tw = bbox[2] - bbox[0]
            x = (800 - tw) // 2
            draw.text((x, y), line, fill=text_color, font=font)
            y += bbox[3] - bbox[1] + 40

        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return send_file(buf, mimetype='image/png', as_attachment=True, download_name='poster.png')
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
